"""Fill the KSP Datathon 2026 Prototype Submission Template with Crime AI content."""
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from pptx.oxml import parse_xml

ROOT = Path(__file__).parent
TEMPLATE = ROOT / "KSP Datathon 2026 _ Prototype Submission Template.pptx"
OUTPUT = ROOT / "KSP_Datathon_2026_CrimeAI_Prototype_Submission.pptx"
ASSETS = ROOT / "assets"

NAVY = RGBColor(15, 42, 78)
SAFFRON = RGBColor(232, 119, 34)
MUTED = RGBColor(55, 70, 95)
WHITE = RGBColor(255, 255, 255)


def set_run_font(run, size=14, bold=False, color=MUTED):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    rPr = run._r.get_or_add_rPr()
    # Ensure East Asian / latin font
    ea = rPr.find(qn("a:latin"))
    if ea is None:
        latin = parse_xml(
            f'<a:latin xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="Segoe UI"/>'
        )
        rPr.insert(0, latin)


def clear_extra_paragraphs(tf, keep=1):
    # remove paragraphs beyond keep
    while len(tf.paragraphs) > keep:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)


def add_textbox(slide, left, top, width, height, lines, default_size=13):
    """lines: list of (text, size, bold, color, level) or plain strings."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, bold, color, level = item, default_size, False, MUTED, 0
        else:
            text = item[0]
            size = item[1] if len(item) > 1 else default_size
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else MUTED
            level = item[4] if len(item) > 4 else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.clear()
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        set_run_font(run, size=size, bold=bold, color=color)
        p.space_after = Pt(4)
    return box


def append_to_existing(shape, extra_lines, size=13, bold=False, color=MUTED, level=0):
    tf = shape.text_frame
    for item in extra_lines:
        if isinstance(item, str):
            text, sz, bd, col, lvl = item, size, bold, color, level
        else:
            text = item[0]
            sz = item[1] if len(item) > 1 else size
            bd = item[2] if len(item) > 2 else bold
            col = item[3] if len(item) > 3 else color
            lvl = item[4] if len(item) > 4 else level
        p = tf.add_paragraph()
        p.level = lvl
        run = p.add_run()
        run.text = text
        set_run_font(run, size=sz, bold=bd, color=col)
        p.space_after = Pt(3)


def fill_label_value(shape, mapping):
    """Replace paragraphs that start with known labels, appending values."""
    tf = shape.text_frame
    for p in tf.paragraphs:
        label = p.text.strip()
        for key, value in mapping.items():
            if label.startswith(key):
                # clear runs and rewrite
                for r in list(p.runs):
                    r.text = ""
                if p.runs:
                    p.runs[0].text = f"{key} {value}"
                    set_run_font(p.runs[0], size=14, bold=True, color=NAVY)
                else:
                    run = p.add_run()
                    run.text = f"{key} {value}"
                    set_run_font(run, size=14, bold=True, color=NAVY)
                break


def add_picture(slide, path, left, top, width):
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def main():
    prs = Presentation(str(TEMPLATE))
    slides = list(prs.slides)

    # -------- Slide 1: Team Details --------
    s1 = slides[0]
    team_shape = s1.shapes[0]
    # Expand the text box upward/taller for readability
    team_shape.top = Inches(2.85)
    team_shape.height = Inches(2.55)
    fill_label_value(
        team_shape,
        {
            "Team name:": "[TEAM NAME]",
            "Team leader name:": "[TEAM LEADER]",
            "Team size:": "[N]",
            "Problem Statement:": "Intelligent Conversational AI for KSP Crime Database",
        },
    )

    # -------- Slide 2: Brief about the solution --------
    s2 = slides[1]
    # Expand title box slightly then add content below
    add_textbox(
        s2,
        0.35,
        1.55,
        9.3,
        3.7,
        [
            ("Crime AI is a conversational intelligence platform for Karnataka SCRB investigators.", 15, True, NAVY),
            ("Instead of static dashboards and manual SQL, officers ask in English or Kannada — by text or voice — and receive explainable answers grounded in the live crime database.", 13, False, MUTED),
            ("", 8),
            ("What it does", 14, True, SAFFRON),
            ("• Turns natural-language questions into safe, read-only SQL / Cypher / analytics queries", 13),
            ("• Surfaces criminal networks, crime trends, hotspots, and early warnings", 13),
            ("• Keeps conversation context across follow-ups (\"show only pending ones\")", 13),
            ("• Exports investigation history as an auditable PDF with role-based access control", 13),
            ("", 8),
            ("Built for SCRB scale: 1100+ police stations, FIR-centric relational schema, graph relationships, bilingual voice.", 13, True, NAVY),
        ],
    )

    # -------- Slide 3: Opportunities --------
    s3 = slides[2]
    # Shrink title block to "Opportunities" only, then add answers
    title_shape = [sh for sh in s3.shapes if sh.has_text_frame][0]
    for i, p in enumerate(list(title_shape.text_frame.paragraphs)):
        for r in p.runs:
            r.text = ""
        if i == 0:
            if p.runs:
                p.runs[0].text = "Opportunities"
                set_run_font(p.runs[0], size=16, bold=True, color=NAVY)
            else:
                run = p.add_run()
                run.text = "Opportunities"
                set_run_font(run, size=16, bold=True, color=NAVY)
    title_shape.height = Inches(0.5)
    add_textbox(
        s3,
        0.3,
        1.45,
        9.4,
        3.9,
        [
            ("How different is it from existing ideas?", 14, True, SAFFRON),
            ("Most crime systems stop at dashboards. Crime AI adds a bilingual voice+chat agent that queries Postgres + Neo4j, retains context, and returns the exact SQL/Cypher used — not a black-box chatbot.", 12),
            ("", 6),
            ("How will it solve the problem?", 14, True, SAFFRON),
            ("Investigators ask questions the way they think (\"robbery FIRs in Hubballi\", \"accused history\", \"criminal network\"). The agent routes to the right tool, returns results in seconds, and supports proactive hotspot / early-warning analysis.", 12),
            ("", 6),
            ("USP", 14, True, SAFFRON),
            ("• Unified stack: NL chat + Kannada/English voice + graph networks + predictive dashboard", 12),
            ("• Explainable AI: every answer ships with tool name, query, and audit trail", 12),
            ("• Role-secure PDF evidence export aligned to police hierarchy (Investigator → Admin)", 12),
            ("• Built on real KSP FIR ER schema (Neon) with Neo4j derived for relationship traversal", 12),
        ],
    )

    # -------- Slide 4: Features --------
    s4 = slides[3]
    add_textbox(
        s4,
        0.25,
        1.65,
        4.6,
        3.6,
        [
            ("Conversation & Voice", 14, True, SAFFRON),
            ("• NL chatbot — English + Kannada", 12),
            ("• Voice-enabled (Sarvam STT/TTS)", 12),
            ("• Context-aware multi-turn chat", 12),
            ("• Streaming token responses", 12),
            ("", 6),
            ("Intelligence", 14, True, SAFFRON),
            ("• Crime pattern discovery", 12),
            ("• Criminal network analysis (Neo4j)", 12),
            ("• Socio-demographic insights", 12),
            ("• Predictive / early-warning signals", 12),
        ],
    )
    add_textbox(
        s4,
        5.1,
        1.65,
        4.6,
        3.6,
        [
            ("Visualization & Evidence", 14, True, SAFFRON),
            ("• Officer dashboard (ECharts)", 12),
            ("• Crime hotspot map (Leaflet)", 12),
            ("• Interactive network graph", 12),
            ("• PDF export of investigation history", 12),
            ("", 6),
            ("Trust & Security", 14, True, SAFFRON),
            ("• Explainable AI with SQL/Cypher audit", 12),
            ("• Role-based access (5 police roles)", 12),
            ("• PDF export gated for SHO and above", 12),
            ("• Read-only DB queries for safety", 12),
        ],
    )

    # -------- Slide 5: Process flow --------
    s5 = slides[4]
    add_picture(s5, ASSETS / "process_flow.png", 0.25, 1.55, 9.5)

    # -------- Slide 6: Wireframes --------
    s6 = slides[5]
    add_picture(s6, ASSETS / "wireframes.png", 0.25, 1.65, 9.5)

    # -------- Slide 7: Architecture --------
    s7 = slides[6]
    add_picture(s7, ASSETS / "architecture.png", 0.25, 1.6, 9.5)

    # -------- Slide 8: Technologies --------
    s8 = slides[7]
    add_textbox(
        s8,
        0.25,
        1.55,
        4.7,
        3.8,
        [
            ("Frontend & Backend", 14, True, SAFFRON),
            ("• Next.js 14 + Tailwind CSS", 12),
            ("• FastAPI (Python 3.11+)", 12),
            ("• Apache ECharts + Leaflet", 12),
            ("", 6),
            ("AI & Voice", 14, True, SAFFRON),
            ("• Gemini 2.5 Flash + LangGraph", 12),
            ("• Sarvam AI — STT (saarika) & TTS (bulbul)", 12),
            ("• LiveKit Agents + Silero VAD", 12),
        ],
    )
    add_textbox(
        s8,
        5.1,
        1.55,
        4.7,
        3.8,
        [
            ("Data & Analytics", 14, True, SAFFRON),
            ("• Neon PostgreSQL + pgvector", 12),
            ("• Neo4j Aura (relationship graph)", 12),
            ("• Pandas / DuckDB analytics", 12),
            ("", 6),
            ("Platform & Auth", 14, True, SAFFRON),
            ("• Zoho Catalyst AppSail + Client Hosting", 12),
            ("• JWT role-based authentication", 12),
            ("• ReportLab PDF generation", 12),
        ],
    )

    # -------- Slide 9: Catalyst services --------
    s9 = slides[8]
    add_textbox(
        s9,
        0.3,
        1.55,
        9.4,
        3.8,
        [
            ("Zoho Catalyst services in this solution", 14, True, SAFFRON),
            ("", 4),
            ("1. AppSail — hosts the FastAPI backend (crime-ai-api) for chat, voice, dashboard, graph, and PDF APIs.", 13),
            ("2. Client Hosting / Web App — serves the Next.js static frontend (chat, dashboard, network, map).", 13),
            ("3. File Store (planned / integration path) — store investigation PDF exports, reports, and conversation archives.", 13),
            ("4. Authentication (integration path) — Catalyst Auth for police roles; prototype currently uses JWT with the same role model (Investigator, SHO, DSP, Analyst, Administrator) so Catalyst Auth can replace the issuer without changing route logic.", 13),
            ("5. Serverless Functions (optional next) — scheduled hotspot jobs, notification hooks, report generation workers.", 13),
            ("", 6),
            ("Deployed API (AppSail): https://crime-ai-api-10130083497.development.catalystappsail.com", 12, True, NAVY),
        ],
    )

    # -------- Slide 10: Cost --------
    s10 = slides[9]
    add_textbox(
        s10,
        0.3,
        1.55,
        9.4,
        3.8,
        [
            ("Estimated monthly operating cost (prototype / pilot scale)", 14, True, SAFFRON),
            ("", 4),
            ("• Zoho Catalyst (AppSail + Client Hosting) - hackathon / free-tier or low paid plan ~ Rs 0-2,000", 12),
            ("• Neon PostgreSQL (free tier -> Launch) ~ Rs 0-1,500", 12),
            ("• Neo4j Aura Free / small paid ~ Rs 0-2,000", 12),
            ("• Google Gemini 2.5 Flash API (interactive chat volume) ~ Rs 1,000-4,000", 12),
            ("• Sarvam AI STT/TTS (voice minutes) ~ Rs 1,000-3,000", 12),
            ("• LiveKit Cloud (realtime voice rooms) ~ Rs 0-2,000", 12),
            ("", 6),
            ("Indicative total for pilot: ~ Rs 3,000 - 15,000 / month", 13, True, NAVY),
            ("Production statewide scale would rise with concurrent users, voice minutes, and DB size; architecture remains the same.", 11),
        ],
    )

    # -------- Slide 11: Snapshots --------
    s11 = slides[10]
    add_picture(s11, ASSETS / "snapshots.png", 0.25, 1.6, 9.5)

    # -------- Slide 12: Performance --------
    s12 = slides[11]
    add_textbox(
        s12,
        0.3,
        1.5,
        9.4,
        3.9,
        [
            ("Prototype performance / benchmarking", 14, True, SAFFRON),
            ("", 4),
            ("Chat latency (measured)", 13, True, NAVY),
            ("• Steady-state chat turn ≈ 5s (down from ~12s) after disabling Gemini extended thinking and adding Neon connection pooling.", 12),
            ("• Each turn: 3 Gemini calls (route → SQL/Cypher → answer synthesis) + one DB/graph query.", 12),
            ("", 4),
            ("Key optimizations", 13, True, NAVY),
            ("• thinking_budget=0 on Gemini 2.5 Flash for routing / schema-grounded SQL / short summaries", 12),
            ("• psycopg connection pool — avoids ~3s cold TLS handshake to Neon per query", 12),
            ("• Streaming SSE (/chat/stream) so investigators see tokens immediately", 12),
            ("", 4),
            ("Demo data scale", 13, True, NAVY),
            ("• ~70+ FIRs across multiple districts / 9 crime types / ~15 months + planted repeat offenders for network demos", 12),
            ("• Voice path: Sarvam STT→LangGraph→TTS verified; LiveKit room path implemented pending cloud credentials", 12),
        ],
    )

    # -------- Slide 13: Links --------
    s13 = slides[12]
    link_shape = [sh for sh in s13.shapes if sh.has_text_frame][0]
    for i, p in enumerate(list(link_shape.text_frame.paragraphs)):
        if i == 0:
            continue
        for r in p.runs:
            r.text = ""
    link_shape.height = Inches(0.5)
    add_textbox(
        s13,
        0.3,
        1.7,
        9.4,
        3.5,
        [
            ("GitHub Public Repository", 14, True, SAFFRON),
            ("[Add public GitHub URL]", 13, False, NAVY),
            ("", 8),
            ("Demo Video Link (3 Minutes)", 14, True, SAFFRON),
            ("[Add YouTube / Drive demo video URL]", 13, False, NAVY),
            ("", 8),
            ("Deployed Link", 14, True, SAFFRON),
            ("API (Catalyst AppSail): https://crime-ai-api-10130083497.development.catalystappsail.com", 12, False, NAVY),
            ("Frontend (Catalyst Client): [Add hosted frontend URL]", 12, False, NAVY),
            ("Local demo: frontend :3000  ·  backend :8000", 12),
        ],
    )

    # -------- Slide 14: Future --------
    s14 = slides[13]
    add_textbox(
        s14,
        0.25,
        1.7,
        9.5,
        3.6,
        [
            ("Additional details / future development", 14, True, SAFFRON),
            ("", 4),
            ("• Swap JWT issuer for Zoho Catalyst Authentication without changing role-gated route logic", 12),
            ("• Persist PDFs and conversation archives to Catalyst File Store", 12),
            ("• Complete LiveKit realtime voice room testing with production credentials", 12),
            ("• Merge multi-step Gemini calls into fewer structured-output round-trips for <3s answers", 12),
            ("• Statewide data ingestion from 1100+ stations + scheduled hotspot / early-warning jobs", 12),
            ("• Deeper behavioral profiling & MO similarity via pgvector over BriefFacts", 12),
            ("• Kannada UI polish and offline-capable field investigator mode", 12),
        ],
    )

    # -------- Slide 15: Thank you / contact --------
    s15 = slides[14]
    blank = [sh for sh in s15.shapes if sh.has_text_frame][0]
    for p in blank.text_frame.paragraphs:
        for r in p.runs:
            r.text = ""
    if blank.text_frame.paragraphs[0].runs:
        blank.text_frame.paragraphs[0].runs[0].text = "Thank You"
        set_run_font(blank.text_frame.paragraphs[0].runs[0], size=28, bold=True, color=NAVY)
    else:
        run = blank.text_frame.paragraphs[0].add_run()
        run.text = "Thank You"
        set_run_font(run, size=28, bold=True, color=NAVY)
    blank.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_textbox(
        s15,
        1.5,
        3.3,
        7.0,
        1.5,
        [
            ("Crime AI — Intelligent Conversational AI for KSP Crime Database", 13, True, NAVY),
            ("KSP Datathon 2026 Prototype Submission", 12, False, MUTED),
            ("Questions welcome — happy to demo live.", 12, False, MUTED),
        ],
    )
    # center the added box paragraphs
    for sh in s15.shapes:
        if sh.has_text_frame and sh != blank and "Crime AI" in (sh.text_frame.text or ""):
            for p in sh.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
